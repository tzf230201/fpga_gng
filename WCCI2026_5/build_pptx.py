# -*- coding: utf-8 -*-
"""
Build an A0 portrait (841 x 1189 mm) PowerPoint poster for:
"Edge-Cell Graph Encoding and Fixed-Point Arithmetic for FPGA Growing Neural Gas"
WCCI 2026.
"""
import os
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))

def mm(v):            # millimetres -> EMU
    return Emu(int(round(v * 36000)))

# ---------- palette ----------
BLUE   = RGBColor(0x0B, 0x3D, 0x6B)   # deep blue (headers / title)
ACCENT = RGBColor(0x1B, 0x7F, 0xB8)   # mid blue
RED    = RGBColor(0xC0, 0x39, 0x2B)   # accent red
GREEN  = RGBColor(0x1E, 0x84, 0x49)
GRAY   = RGBColor(0xEC, 0xF1, 0xF5)   # light panel bg
CODEBG = RGBColor(0xF7, 0xFA, 0xFC)   # formula box bg
DARK   = RGBColor(0x17, 0x25, 0x2E)   # body text
MUTE   = RGBColor(0x5A, 0x6B, 0x76)   # muted gray text
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

MONO = "Consolas"
SANS = "Calibri"

# ---------- presentation ----------
prs = Presentation()
prs.slide_width  = mm(841)
prs.slide_height = mm(1189)
slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank
shapes = slide.shapes

def no_shadow(sh):
    sh.shadow.inherit = False

# ---------- white background ----------
bg = shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, mm(841), mm(1189))
bg.fill.solid(); bg.fill.fore_color.rgb = WHITE
bg.line.fill.background(); no_shadow(bg)

def add_rect(x, y, w, h, fill=None, line=None, line_w=1.0, rounded=False):
    assert w > 0 and h > 0, "non-positive shape size: w=%s h=%s" % (w, h)
    shp = shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        mm(x), mm(y), mm(w), mm(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    no_shadow(shp)
    if rounded:
        try:
            shp.adjustments[0] = 0.04
        except Exception:
            pass
    return shp

def set_runs(p, segs, size, default_color=DARK, font=SANS):
    for (txt, bold, col) in segs:
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = font
        r.font.color.rgb = col if col else default_color

def add_text(x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = shapes.add_shape(MSO_SHAPE.RECTANGLE, mm(x), mm(y), mm(w), mm(h))
    tb.fill.background(); tb.line.fill.background(); no_shadow(tb)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for side in ("left", "right", "top", "bottom"):
        setattr(tf, "margin_" + side, mm(3))
    return tf

def para(tf, segs, size=23, bullet=False, align=PP_ALIGN.LEFT,
         space_before=2, space_after=7, color=DARK, first=False, line=1.03,
         font=SANS):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    try:
        p.line_spacing = line
    except Exception:
        pass
    if bullet:
        segs = [("▸  ", True, ACCENT)] + list(segs)
    set_runs(p, segs, size, default_color=color, font=font)
    return p

# ---------- block (panel + header) ----------
def block(x, y, w, h, title, body_top_pad=26):
    add_rect(x, y, w, h, fill=GRAY, line=BLUE, line_w=2.4, rounded=True)
    hdr = add_rect(x, y, w, 27, fill=BLUE, rounded=True)
    htf = hdr.text_frame
    htf.word_wrap = True
    htf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = htf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    set_runs(p, [(title, True, WHITE)], 31)
    return add_text(x + 7, y + body_top_pad, w - 14, h - body_top_pad - 7)

def picture(path, x, y, w, h=None):
    im = Image.open(os.path.join(HERE, path))
    if h is None:
        h = w * im.height / im.width
    shapes.add_picture(os.path.join(HERE, path), mm(x), mm(y), mm(w), mm(h))
    return h

# =====================================================================
# TITLE BANNER
# =====================================================================
add_rect(0, 0, 841, 172, fill=BLUE)
add_rect(0, 172, 841, 6, fill=ACCENT)

ttf = add_text(30, 12, 781, 92, anchor=MSO_ANCHOR.MIDDLE)
para(ttf, [("Edge-Cell Graph Encoding and Fixed-Point Arithmetic", True, WHITE)],
     size=60, align=PP_ALIGN.CENTER, space_after=4, first=True)
para(ttf, [("for FPGA Growing Neural Gas", True, WHITE)],
     size=60, align=PP_ALIGN.CENTER, space_before=0)

atf = add_text(30, 106, 781, 30, anchor=MSO_ANCHOR.MIDDLE)
para(atf, [("Teuku Zikri Fatahillah", True, WHITE), ("¹   ", False, ACCENT),
           ("Raditya Artha Rochmanto", True, WHITE), ("¹²   ", False, ACCENT),
           ("Achmad Fahrul Aji", True, WHITE), ("¹²   ", False, ACCENT),
           ("Anhar Risnumawan", True, WHITE), ("¹³   ", False, ACCENT),
           ("Chyan Zheng Siow", True, WHITE), ("¹⁴   ", False, ACCENT),
           ("Naoyuki Kubota", True, WHITE), ("¹", False, ACCENT)],
     size=27, align=PP_ALIGN.CENTER, first=True, space_after=0)

itf = add_text(30, 138, 781, 30, anchor=MSO_ANCHOR.MIDDLE)
para(itf, [("¹ Tokyo Metropolitan University, Japan      "
            "² Politeknik Negeri Semarang, Indonesia      "
            "³ Politeknik Elektronika Negeri Surabaya, Indonesia      "
            "⁴ Changsha Cultural and Creative Arts Vocational College, China",
            False, RGBColor(0xCF, 0xE2, 0xF0))],
     size=19, align=PP_ALIGN.CENTER, first=True, space_after=0)

# =====================================================================
# COLUMN GEOMETRY
# =====================================================================
LX, RX = 20, 431
CW = 390
TOP = 192
BOT = 1174
gap = 12

# ---------------------------------------------------------------------
# LEFT COLUMN
# ---------------------------------------------------------------------
y = TOP

# --- Motivation & Problem  (3 bullets + large figure) ---
h = 300
tf = block(LX, y, CW, h, "Motivation & Problem")
para(tf, [("Edge robots & micro-UAVs need ", False, None),
          ("on-board unsupervised learning", True, BLUE),
          (" — no cloud, no labels.", False, None)],
     size=23, bullet=True, first=True, space_after=8)
para(tf, [("GNG", True, BLUE),
          (" fits (incremental, self-organizing) but is ", False, None),
          ("memory-heavy", True, RED), (".", False, None)],
     size=23, bullet=True, space_after=8)
para(tf, [("Dense adjacency ", False, None), ("O(N²) = 1,600 B", True, RED),
          (" and floating-point are costly on small FPGAs.", False, None)],
     size=23, bullet=True, space_after=4)

# --- accurate edge-cell concept figure (replaces inconsistent problem.png) ---
gtitle = add_text(LX + 7, y + 116, CW - 14, 14)
para(gtitle, [("Our fix — store only the upper-triangular half:", True, BLUE)],
     size=20, first=True, align=PP_ALIGN.CENTER, space_after=0)

def draw_edgecell_grid(gx, gy, n=7, cell=15.5, gap=1.6):
    light = RGBColor(0xD9, 0xDF, 0xE3)
    for i in range(n):          # row
        for j in range(n):      # col
            cx = gx + j * (cell + gap)
            cy = gy + i * (cell + gap)
            if j > i:
                fill = ACCENT          # upper triangle: stored (edge-cell)
            elif j == i:
                fill = MUTE            # diagonal: no self-edge
            else:
                fill = light          # lower triangle: not stored (symmetric)
            add_rect(cx, cy, cell, cell, fill=fill,
                     line=WHITE, line_w=1.0)

GN, GC, GG = 7, 15.5, 1.6
gsize = GN * GC + (GN - 1) * GG
gx = LX + 24
gy = y + 134
draw_edgecell_grid(gx, gy, GN, GC, GG)

# legend + numbers to the right of the grid
lx = gx + gsize + 16
add_rect(lx, gy + 4, 13, 13, fill=ACCENT, line=WHITE, line_w=1.0)
lg = add_text(lx + 16, gy - 1, CW - (lx - LX) - 10, 24)
para(lg, [("stored: edge-cell (age+1)", False, DARK)], size=19, first=True, space_after=2)
add_rect(lx, gy + 32, 13, 13, fill=RGBColor(0xD9, 0xDF, 0xE3), line=WHITE, line_w=1.0)
lg2 = add_text(lx + 16, gy + 27, CW - (lx - LX) - 10, 24)
para(lg2, [("not stored (symmetric half)", False, DARK)], size=19, first=True, space_after=2)
lg3 = add_text(lx - 2, gy + 64, CW - (lx - LX) - 6, 50)
para(lg3, [("Dense ", False, None), ("O(N²)", True, RED),
           (" = 1,600 B", False, None)], size=19, first=True, space_after=3)
para(lg3, [("→  ", True, GREEN), ("780 B", True, GREEN),
           ("  half-adjacency", False, None)], size=19, space_after=3)
para(lg3, [("= ", False, None), ("−51% edge storage", True, GREEN)], size=20)
y += h + gap

# --- Contributions ---
h = 215
tf = block(LX, y, CW, h, "Our Contributions")
para(tf, [("A fully integer, single-FSM GNG engine", True, BLUE),
          (" in synthesizable VHDL on the ", False, None),
          ("Sipeed Tang Nano 9K", True, BLUE),
          (" — no soft-core CPU, no FPU.", False, None)],
     size=22, first=True, space_after=8)
para(tf, [("Edge-cell encoding", True, BLUE),
          (": half-adjacency as 8-bit age+1 — ", False, None),
          ("780 B vs 1,600 B.", True, RED)], size=22, bullet=True)
para(tf, [("Shift-based Q8/Q16", True, BLUE),
          (" rates ⇒ single-cycle DSP multiplies.", False, None)],
     size=22, bullet=True)
para(tf, [("Fused update pass", True, BLUE),
          (": decay + aging + move + prune in one scan.", False, None)],
     size=22, bullet=True)
y += h + gap

# --- GNG Algorithm Core ---
h = 200
tf = block(LX, y, CW, h, "GNG Algorithm Core")
para(tf, [("Graph ", False, None), ("G = (V, E)", True, BLUE),
          (", weight vectors wᵢ ∈ ℝ². Per input ξ:", False, None)],
     size=22, first=True, space_after=6)
for seg in [
    [("Find winner ", False, None), ("s₁", True, BLUE),
     (" and runner-up ", False, None), ("s₂", True, BLUE), (".", False, None)],
    [("Accumulate / decay error; move s₁ toward ξ.", False, None)],
    [("Age edges of s₁; move neighbors; ", False, None),
     ("prune over-age.", True, RED)],
    [("Connect (s₁, s₂); every λ steps insert a node.", False, None)],
]:
    para(tf, seg, size=22, bullet=True, space_after=4)
y += h + gap

# --- FPGA Architecture (FSM pipeline) ---
h = BOT - y
tf = block(LX, y, CW, h, "FPGA Architecture")
para(tf, [("Single synchronous ", False, None), ("FSM datapath @ 27 MHz", True, BLUE),
          (" (no PLL). Three BRAMs: node / edge / dataset.  ", False, None),
          ("80-bit node word", True, BLUE),
          (" = x · y · act · deg · err.", False, None)],
     size=22, first=True, space_after=6)

phases = ["INIT", "SAMPLE", "WIN_SCAN", "UPDATE",
          "NB_SCAN", "CONNECT", "INSERT", "DBG/SNAP"]
cw, ch = 88, 36
cgx, cgy = 8, 26
cx0 = LX + (CW - (4 * cw + 3 * cgx)) / 2
cy0 = y + 100
for i, name in enumerate(phases):
    row, col = divmod(i, 4)
    cx = cx0 + col * (cw + cgx)
    cy = cy0 + row * (ch + cgy)
    chip = add_rect(cx, cy, cw, ch, fill=WHITE, line=ACCENT, line_w=2.2, rounded=True)
    ctf = chip.text_frame; ctf.word_wrap = True
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for s in ("left", "right", "top", "bottom"):
        setattr(ctf, "margin_" + s, mm(1))
    pp = ctf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
    set_runs(pp, [(name, True, BLUE)], 18)
    if col < 3:
        ar = shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, mm(cx + cw + 0.5),
                              mm(cy + ch/2 - 4.5), mm(cgx - 1), mm(9))
        ar.fill.solid(); ar.fill.fore_color.rgb = ACCENT
        ar.line.fill.background(); no_shadow(ar)
wa = shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                      mm(cx0 + 3*(cw+cgx) + cw/2 - 4.5), mm(cy0 + ch + 3),
                      mm(9), mm(cgy - 6))
wa.fill.solid(); wa.fill.fore_color.rgb = ACCENT
wa.line.fill.background(); no_shadow(wa)

ftf = add_text(LX + 7, cy0 + 2*ch + cgy + 8, CW - 14, 34)
para(ftf, [("Deterministic, bounded per-sample latency — "
            "each BRAM read-modify-write ≥ 3 cycles.", False, None)],
     size=20, first=True, color=DARK)

# ---------------------------------------------------------------------
# RIGHT COLUMN
# ---------------------------------------------------------------------
y = TOP

# --- Key Idea: Edge-Cell Encoding ---
h = 198
tf = block(RX, y, CW, h, "Key Idea: Edge-Cell Encoding")
para(tf, [("Store only the ", False, None), ("upper triangle", True, BLUE),
          (" of the adjacency, one byte per edge, ", False, None),
          ("age+1", True, BLUE), (" encoded:", False, None)],
     size=22, first=True, space_after=6)
# highlighted result line
add_rect(RX + 10, y + 70, CW - 20, 30, fill=WHITE, line=ACCENT, line_w=2.0, rounded=True)
rb = add_text(RX + 10, y + 70, CW - 20, 30, anchor=MSO_ANCHOR.MIDDLE)
para(rb, [("E = N(N−1)/2 = ", False, DARK),
          ("780 B", True, RED),
          ("   (vs 1,600 B dense)", False, MUTE)],
     size=24, font=MONO, align=PP_ALIGN.CENTER, first=True, space_after=0)
tf2 = add_text(RX + 7, y + 106, CW - 14, h - 112)
para(tf2, [("cell = 0", True, BLUE), (" no edge   ·   ", False, None),
           ("cell = v>0", True, BLUE), (" active, age = v−1", False, None)],
     size=21, bullet=True, first=True)
para(tf2, [("No active bit; over-age prune = one compare ", False, None),
           ("(v > 51)", True, BLUE), (".", False, None)], size=21, bullet=True)
para(tf2, [("Fused NB_SCAN", True, RED),
           (" folds decay + aging + move + prune in one scan.", False, None)],
     size=21, bullet=True)
y += h + gap

# --- Integer Fixed-Point Arithmetic (clean formula box) ---
h = 200
tf = block(RX, y, CW, h, "Integer Fixed-Point Arithmetic")
para(tf, [("16-bit signed coords in ", False, None), ("[−1000, 1000]", True, BLUE),
          (" — shift-scaled rates, no FPU and no divider:", False, None)],
     size=22, first=True, space_after=4)

# formula panel
fx, fy, fw, fh = RX + 10, y + 66, CW - 20, 96
add_rect(fx, fy, fw, fh, fill=CODEBG, line=ACCENT, line_w=2.0, rounded=True)
fbox = add_text(fx + 4, fy + 3, fw - 8, fh - 6, anchor=MSO_ANCHOR.MIDDLE)
def fline(label, expr, note, first=False):
    p = fbox.paragraphs[0] if first and not fbox.paragraphs[0].runs else fbox.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.space_before = Pt(2); p.space_after = Pt(2)
    try: p.line_spacing = 1.05
    except Exception: pass
    set_runs(p, [(label, True, BLUE), (expr, False, DARK), (note, False, GREEN)],
             20, font=MONO)
fline("dist     ", "d² = Δx² + Δy²", "    9×9 DSP, 35-bit", first=True)
fline("winner   ", "Δ = (77·δ) >> 8", "   ε ≈ 0.301  Q8")
fline("neighbor ", "Δ = (66·δ) >> 16", "  ε ≈ 0.001  Q16")
fline("decay    ", "err −= err >> 8", "   β ≈ 0.996")
y += h + gap

# --- Results ---
h = 234
block(RX, y, CW, h, "Results")

def style_table(tbl, font_size, col_aligns=None):
    for ri, row in enumerate(tbl.rows):
        row.height = mm(12)
        for ci, cell in enumerate(row.cells):
            cell.margin_left = mm(2.5); cell.margin_right = mm(2.5)
            cell.margin_top = mm(0.5); cell.margin_bottom = mm(0.5)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            if col_aligns:
                p.alignment = col_aligns[ci]
            for r in p.runs:
                r.font.size = Pt(font_size); r.font.name = SANS
                if ri == 0:
                    r.font.bold = True; r.font.color.rgb = WHITE
                else:
                    r.font.color.rgb = DARK
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = BLUE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if ri % 2 else RGBColor(0xDD, 0xE8, 0xF1)

def make_table(x, y, w, data, col_w, font_size=18, aligns=None):
    rows, cols = len(data), len(data[0])
    gt = shapes.add_table(rows, cols, mm(x), mm(y), mm(w), mm(12*rows)).table
    gt.first_row = False; gt.horz_banding = False
    for ci, cwid in enumerate(col_w):
        gt.columns[ci].width = mm(cwid)
    for ri, rowd in enumerate(data):
        for ci, val in enumerate(rowd):
            gt.cell(ri, ci).text = val
    style_table(gt, font_size, col_aligns=aligns)
    return gt

L = PP_ALIGN.LEFT; C = PP_ALIGN.CENTER; R = PP_ALIGN.RIGHT
cap = add_text(RX + 8, y + 30, 182, 12)
para(cap, [("Memory footprint (Nₘₐₓ=40)", True, BLUE)], size=18, first=True,
     align=PP_ALIGN.CENTER, space_after=0)
cap2 = add_text(RX + 200, y + 30, 182, 12)
para(cap2, [("FPGA utilization (GW1NR-9C)", True, BLUE)], size=18, first=True,
     align=PP_ALIGN.CENTER, space_after=0)
mem = [["", "Edge", "Total"],
       ["Dense adj.", "1600 B", "2000 B"],
       ["Edge-cell", "780 B", "1180 B"],
       ["Reduction", "51%", "41%"]]
make_table(RX + 8, y + 44, 182, mem, [70, 56, 56], font_size=18, aligns=[L, R, R])
util = [["Resource", "Used", "%"],
        ["LUT", "3205", "38%"],
        ["FF", "1408", "22%"],
        ["BSRAM", "6", "24%"],
        ["DSP", "5", "50%"]]
make_table(RX + 200, y + 44, 182, util, [96, 50, 36], font_size=18, aligns=[L, R, R])

rcap = add_text(RX + 8, y + 112, 374, 12)
para(rcap, [("Topology quality & runtime (27 MHz)", True, BLUE)], size=18,
     first=True, align=PP_ALIGN.CENTER, space_after=0)
q = [["Dataset", "Nodes", "Edges", "QE", "TE", "µs/sample"],
     ["Two Moons", "40", "37", "0.146", "0.00", "624"],
     ["Concentric Circles", "40", "40", "0.200", "0.04", "627"]]
make_table(RX + 8, y + 126, 374, q, [120, 50, 50, 50, 46, 58], font_size=18,
           aligns=[L, C, C, C, C, C])
note = add_text(RX + 8, y + 170, 374, 24)
para(note, [("Deterministic ~625 µs/sample  ⇒  ~1,600 samples/s "
             "on a sub-10K-LUT FPGA.", True, GREEN)], size=20, first=True,
     align=PP_ALIGN.CENTER)
y += h + gap

# --- Learned Topologies (larger plots) ---
h = 205
block(RX, y, CW, h, "Learned Topologies")
iw = 184
ix0 = RX + (CW - (2*iw + 8)) / 2
ih = picture("two_moons_result.png", ix0, y + 34, iw)
picture("concentric_circles_result.png", ix0 + iw + 8, y + 34, iw)
icap = add_text(RX + 8, y + 36 + ih, CW - 16, 22)
para(icap, [("Two Moons (left) and Concentric Circles (right), ", False, None),
            ("Nₘₐₓ = 40", True, BLUE), (".", False, None)],
     size=18, first=True, align=PP_ALIGN.CENTER)
y += h + gap

# --- Conclusion + QR ---
h = BOT - y
tf = block(RX, y, CW - 78, h, "Conclusion")
para(tf, [("A ", False, None),
          ("fully integer, sub-10K-LUT FPGA GNG engine", True, BLUE),
          (". Edge-cell encoding ", False, None),
          ("halves edge storage", True, RED),
          ("; Q8/Q16 fixed-point keeps QE/TE low with deterministic latency — "
           "a foundation for ", False, None),
          ("continual clustering at the edge.", True, BLUE)],
     size=21, first=True, space_after=5)
para(tf, [("Supported by JST Moonshot R&D, grant JPMJMS2034.", False, MUTE)],
     size=16)

# QR block (right of conclusion)
qx, qy, qsz = RX + CW - 72, y, 72
add_rect(qx, qy, qsz, h, fill=WHITE, line=BLUE, line_w=2.4, rounded=True)
qpic = 56
picture("qr_github.png", qx + (qsz - qpic)/2, qy + 8, qpic, qpic)
qlab = add_text(qx + 3, qy + 8 + qpic + 1, qsz - 6, h - qpic - 14, anchor=MSO_ANCHOR.TOP)
para(qlab, [("Source code", True, BLUE)], size=15, first=True,
     align=PP_ALIGN.CENTER, space_after=0)
para(qlab, [("& datasets", True, BLUE)], size=15,
     align=PP_ALIGN.CENTER, space_after=0, space_before=0)

# =====================================================================
out = os.path.join(HERE, "poster.pptx")
prs.save(out)
print("saved:", out)
